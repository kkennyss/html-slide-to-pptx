"""
screenshot_to_pptx.py - Convert HTML slides to PPTX via browser screenshots.

Usage:
    python screenshot_to_pptx.py <html_url> [--output <path>] [--port <chrome_port>]

Requires:
    - Chrome/Chromium installed at C:\Program Files\Google\Chrome\Application\chrome.exe
    - pip packages: requests, websocket-client, python-pptx
    - An HTTP server serving the HTML (this script does NOT start one)
"""
import subprocess, time, json, os, base64, requests, sys, argparse, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

CHROME_PATH = 'C:/Program Files/Google/Chrome/Application/chrome.exe'

def screenshot_to_pptx(html_url, output_path, chrome_port=9225):
    """Screenshot each .slide div from HTML and create PPTX."""

    OUT_DIR = os.path.join(os.path.dirname(output_path) or '.', '_ppt_slides')
    os.makedirs(OUT_DIR, exist_ok=True)

    # ---- 1. Launch headless Chrome ----
    proc = subprocess.Popen([
        CHROME_PATH, '--headless=new',
        f'--remote-debugging-port={chrome_port}',
        '--remote-allow-origins=*', '--no-sandbox',
        '--disable-gpu', '--hide-scrollbars',
        html_url
    ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    time.sleep(5)

    # ---- 2. Connect to Chrome DevTools ----
    r = requests.get(f'http://127.0.0.1:{chrome_port}/json', timeout=5)
    pages = r.json()
    # Find the HTML page tab
    target = None
    for p in pages:
        if html_url in p.get('url', ''):
            target = p; break
    if not target:
        target = pages[-1]  # fallback to newest tab

    import websocket
    ws = websocket.create_connection(target['webSocketDebuggerUrl'], timeout=30)

    def cdp(method, params=None):
        if params is None: params = {}
        cid = int(time.time() * 1000) % 100000
        ws.send(json.dumps({"id": cid, "method": method, "params": params}))
        while True:
            resp = json.loads(ws.recv())
            if resp.get("id") == cid:
                if "error" in resp:
                    print(f"  CDP Error [{method}]: {resp['error']['message']}")
                return resp.get("result")

    time.sleep(3)

    # ---- 3. Debug: check if page loaded ----
    r = cdp("Runtime.evaluate", {"expression": "document.title", "returnByValue": True})
    print(f"  Page title: {r.get('result',{}).get('value','n/a') if r else 'n/a'}")

    # ---- 4. Get slide count ----
    r = cdp("Runtime.evaluate", {
        "expression": "document.querySelectorAll('.slide').length",
        "returnByValue": True
    })
    slide_count = r.get("result", {}).get("value", 0) if r else 0
    print(f"  Slide count: {slide_count}")

    if slide_count == 0:
        # Try harder - check body
        r = cdp("Runtime.evaluate", {"expression": "document.body ? document.body.innerHTML.substring(0,300) : 'no body'", "returnByValue": True})
        body = r.get("result",{}).get("value","n/a") if r else "n/a"
        print(f"  Body preview: {body[:200]}...")
        raise RuntimeError("No .slide elements found in HTML")

    r = cdp("Runtime.evaluate", {
        "expression": """
        JSON.stringify(Array.from(document.querySelectorAll('.slide'), function(el, i) {
            var rect = el.getBoundingClientRect();
            return {
                index: i + 1,
                top: Math.round(rect.top + window.scrollY),
                left: Math.round(rect.left),
                width: Math.round(rect.width),
                height: Math.round(rect.height)
            };
        }))
        """,
        "returnByValue": True
    })
    slides = json.loads(r["result"]["value"])
    print(f"  Found {len(slides)} slides with positions")

    # ---- 5. Hide UI overlays before screenshotting ----
    cdp("Runtime.evaluate", {
        "expression": """
        var els = document.querySelectorAll('#download-btn, .screenshot-hide, [data-hide-screenshot]');
        for (var i = 0; i < els.length; i++) {
            els[i].style.display = 'none';
        }
        """,
        "returnByValue": True
    })
    time.sleep(0.3)

    # ---- 6. Screenshot each slide ----
    for s in slides:
        idx = s['index']
        print(f"  Slide {idx} ({s['width']}x{s['height']})", end='', flush=True)

        cdp("Runtime.evaluate", {
            "expression": f"window.scrollTo(0, {s['top'] - 5})",
            "returnByValue": True
        })
        time.sleep(0.3)

        cdp("Emulation.setDeviceMetricsOverride", {
            "width": s['width'], "height": s['height'] + 5,
            "deviceScaleFactor": 2, "mobile": False
        })
        time.sleep(0.2)

        r = cdp("Page.captureScreenshot", {"format": "png", "fromSurface": True})
        if r and "data" in r:
            img = base64.b64decode(r["data"])
            out = os.path.join(OUT_DIR, f"slide_{idx:02d}.png")
            with open(out, "wb") as f:
                f.write(img)
            print(f"  {len(img)//1024}KB", flush=True)

        cdp("Emulation.clearDeviceMetricsOverride")
        time.sleep(0.15)

    ws.close()
    proc.terminate()
    print(f"\nScreenshots saved to: {OUT_DIR}")

    # ---- 7. Assemble PPTX ----
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9
    prs.slide_height = Emu(6858000)

    images = sorted(
        [f for f in os.listdir(OUT_DIR) if f.endswith('.png')],
        key=lambda x: int(x.split('_')[1].split('.')[0])
    )

    blank = prs.slide_layouts[6]
    for img_name in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            os.path.join(OUT_DIR, img_name),
            Emu(0), Emu(0), prs.slide_width, prs.slide_height
        )

    prs.save(output_path)
    print(f"PPTX saved: {output_path} ({len(images)} slides)")

    # ---- 8. Cleanup temp screenshots ----
    import shutil
    shutil.rmtree(OUT_DIR, ignore_errors=True)

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='HTML slides to PPTX')
    parser.add_argument('html_url', help='URL of the HTML presentation')
    parser.add_argument('--output', '-o', default=None,
                        help='Output PPTX path')
    parser.add_argument('--port', '-p', type=int, default=9225,
                        help='Chrome DevTools port (default: 9225)')
    args = parser.parse_args()

    if not args.output:
        args.output = os.path.expanduser(r'~\Desktop\presentation.pptx')

    screenshot_to_pptx(args.html_url, args.output, args.port)
